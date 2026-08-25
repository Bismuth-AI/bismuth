"""Text extraction for presentation and spreadsheet formats."""

from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path

from bismuth.adapters.parsers.registry import build_extraction, require
from bismuth.domain.document import Extraction, Section
from bismuth.domain.errors import ParserUnavailableError


class PptxParser:
    """Slide decks; one section per slide, with speaker notes attached."""

    @property
    def name(self) -> str:
        return "python-pptx"

    @property
    def extensions(self) -> frozenset[str]:
        return frozenset({".pptx"})

    def warm(self) -> None:
        require("pptx", "Reading .pptx needs python-pptx: pip install 'bismuth-kb[parsers]'")

    def parse(self, path: Path, *, max_chars: int) -> Extraction:
        self.warm()
        from pptx import Presentation

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise ParserUnavailableError(f"{path.name} is not a readable .pptx: {exc}") from exc

        return build_extraction(_slides(presentation), parser=self.name, max_chars=max_chars)


def _slides(presentation: object) -> Iterator[Section]:
    for index, slide in enumerate(presentation.slides):  # type: ignore[attr-defined]
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and (text := shape.text_frame.text.strip()):
                lines.append(text)

        if (
            slide.has_notes_slide
            and slide.notes_slide.notes_text_frame is not None
            and (notes := slide.notes_slide.notes_text_frame.text.strip())
        ):
            lines.append(f"\n[speaker notes] {notes}")

        if lines:
            yield Section(
                heading=f"Slide {index + 1}", text="\n".join(lines), page=index + 1, order=index
            )


class XlsxParser:
    """Extract each worksheet as a Markdown table using cached formula values."""

    @property
    def name(self) -> str:
        return "openpyxl"

    @property
    def extensions(self) -> frozenset[str]:
        return frozenset({".xlsx", ".xlsm"})

    def warm(self) -> None:
        require("openpyxl", "Reading .xlsx needs openpyxl: pip install 'bismuth-kb[parsers]'")

    def parse(self, path: Path, *, max_chars: int) -> Extraction:
        self.warm()
        import openpyxl

        try:
            workbook = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        except Exception as exc:
            raise ParserUnavailableError(
                f"{path.name} is not a readable spreadsheet: {exc}"
            ) from exc

        try:
            return build_extraction(_sheets(workbook), parser=self.name, max_chars=max_chars)
        finally:
            workbook.close()


def _sheets(workbook: object) -> Iterator[Section]:
    for order, sheet in enumerate(workbook.worksheets):  # type: ignore[attr-defined]
        rows: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v).replace("|", "\\|") for v in row]
            if any(cell.strip() for cell in cells):
                rows.append(cells)
        if rows:
            yield Section(heading=str(sheet.title), text=_as_markdown_table(rows), order=order)


def _as_markdown_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    padded = [row + [""] * (width - len(row)) for row in rows]
    header, *body = padded
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
        *("| " + " | ".join(row) + " |" for row in body),
    ]
    return "\n".join(lines)
